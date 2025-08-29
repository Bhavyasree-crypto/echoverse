import streamlit as st
from gtts import gTTS
import os
import docx
import PyPDF2
from pptx import Presentation
from openai import OpenAI

# Initialize OpenAI
client = OpenAI(api_key="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")

# ---------------------- Helper Functions ----------------------

def read_docx(file):
    doc = docx.Document(file)
    return " ".join([para.text for para in doc.paragraphs])

def read_pdf(file):
    reader = PyPDF2.PdfReader(file)
    return " ".join([page.extract_text() for page in reader.pages if page.extract_text()])

def read_pptx(file):
    prs = Presentation(file)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return " ".join(text)

def generate_ai_answer(question):
    response = client.chat.completions.create(
        model="gpt-4o-mini",  # cheaper + fast
        messages=[{"role": "user", "content": question}],
        max_tokens=200
    )
    return response.choices[0].message.content

def text_to_speech(text, lang):
    # gTTS supported langs: 'en' (English), 'hi' (Hindi), 'te' (Telugu)
    tts = gTTS(text=text, lang=lang)
    audio_path = "output.mp3"
    tts.save(audio_path)
    return audio_path

# ---------------------- Streamlit UI ----------------------

st.title("🎧 EchoVerse - AI Audiobook Tool")

choice = st.radio("Choose Input Type:", ["Text", "Question", "Upload File"])

user_text = ""

if choice == "Text":
    user_text = st.text_area("Enter your text here:")

elif choice == "Question":
    question = st.text_input("Ask me a question:")
    if question:
        user_text = generate_ai_answer(question)

elif choice == "Upload File":
    file = st.file_uploader("Upload a document", type=["txt", "docx", "pdf", "pptx"])
    if file:
        if file.type == "text/plain":
            user_text = file.read().decode("utf-8")
        elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            user_text = read_docx(file)
        elif file.type == "application/pdf":
            user_text = read_pdf(file)
        elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            user_text = read_pptx(file)

# ---------------------- Language Selection ----------------------

lang = st.selectbox("Select Language", ["en (English)", "hi (Hindi)", "te (Telugu)"])
lang_code = lang.split(" ")[0]  # pick code (en, hi, te)

# ---------------------- Generate Speech ----------------------

if st.button("Convert to Audio"):
    if user_text:
        audio_path = text_to_speech(user_text, lang_code)
        st.audio(audio_path)
        st.success("✅ Audio generated successfully!")
    else:
        st.warning("Please provide some text, ask a question, or upload a file.")
